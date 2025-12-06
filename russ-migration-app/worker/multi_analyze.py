import subprocess
import json
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from openpyxl import load_workbook
from vm_recommender import recommend_vm_shape  # ✅ import recommender

# ---------------------------------------------------------------------
BASE = Path(__file__).resolve().parent
UPLOADS = BASE.parent / "uploads"
OUTPUTS = Path(__file__).resolve().parent.parent / "outputs"
RUN_PIPELINE = BASE / "run_pipeline.py"
TEMPLATE_PPT = BASE.parent / "analysis_templates" / "template.pptx"

OUTPUTS.mkdir(exist_ok=True)

# ---------------------------------------------------------------------
def analyze_all(cloud="azure"):
    """
    Run pipeline on all uploaded AWR HTML reports and/or manual data.
    Produces summary with metrics + VM recommendations (AWS or Azure).
    """
    print("===============================================================")
    print(f"🧩 Starting Multi-Analysis for all uploads [{cloud.upper()}]")
    print("===============================================================")

    results = []
    awr_files = sorted(list(UPLOADS.glob("*.html")) + list(UPLOADS.glob("*.htm")))

    # ---------------------------------------------------------------------
    # 📥 Check for manual_inputs.json (added for Phase 1A)
    # ---------------------------------------------------------------------
    manual_json = UPLOADS / "manual_inputs.json"
    manual_data = None
    if manual_json.exists():
        try:
            manual_data = json.loads(manual_json.read_text())
            print(f"📘 Found manual_inputs.json with data: {manual_data}")
        except Exception as e:
            print(f"⚠️ Failed to read manual_inputs.json: {e}")

    # ---------------------------------------------------------------------
    # 📊 Handle AWR HTML Files
    # ---------------------------------------------------------------------
    if awr_files:
        for awr in awr_files:
            print("---------------------------------------------------------------")
            print(f"⚙️ Processing file: {awr.name}")

            try:
                # 🧠 Run the single report pipeline
                subprocess.run(
                    ["python3", str(RUN_PIPELINE), "--awr-file", str(awr)],
                    check=True
                )

                # Find the latest generated Excel file
                generated_files = list(OUTPUTS.glob("final_analysis_*.xlsx"))
                if not generated_files:
                    print("⚠️ No generated Excel files found.")
                    continue

                latest_file = max(generated_files, key=lambda f: f.stat().st_mtime)
                wb = load_workbook(latest_file, data_only=True)

                if "AWRData" not in wb.sheetnames:
                    print(f"⚠️ Missing 'AWRData' in {latest_file.name}, skipping.")
                    continue

                ws = wb["AWRData"]

                def val(cell):
                    v = ws[cell].value
                    try:
                        return float(v)
                    except (TypeError, ValueError):
                        return 0.0

                db_name = str(ws["A2"].value or "UNKNOWN")

                # Derived metrics
                memory_gb = (val("I2") + val("J2")) / 1024
                total_iops = val("M2") + val("N2")
                total_throughput = val("K2") + val("L2")
                estimated_vcpus = ((val("E2") / val("D2")) * 1 * 2 * 1.1) / 2 if val("D2") else 0

                print(f"   → Memory: {memory_gb:.2f} GB, IOPS: {total_iops:.2f}, MB/s: {total_throughput:.2f}, vCPUs: {estimated_vcpus:.2f}")

                # Recommend VM
                vm = recommend_vm_shape(cloud, estimated_vcpus, memory_gb)

                result = {
                    "Source": "AWR Report",
                    "Cloud": cloud.upper(),
                    "DB Name": db_name,
                    "Memory (GB)": round(memory_gb, 2),
                    "Total IOPS": round(total_iops, 2),
                    "Throughput (MB/s)": round(total_throughput, 2),
                    "Estimated vCPUs": round(estimated_vcpus, 2),
                    "Recommended VM": vm.get("name", "N/A"),
                    "VM vCPUs": vm.get("vcpus", "N/A"),
                    "VM Memory (GB)": vm.get("memory", "N/A"),
                    "Category": vm.get("category", "N/A"),
                    "Hourly Price (USD)": vm.get("price_per_hour", "N/A"),
                    "Monthly Cost (USD)": vm.get("monthly_cost", "N/A")
                }

                results.append(result)
                print(f"✅ {db_name} → {result['Recommended VM']} @ ${result['Hourly Price (USD)']}/hr")

            except subprocess.CalledProcessError as e:
                print(f"❌ Error processing {awr.name}: {e}")
            except Exception as ex:
                print(f"⚠️ Unexpected error on {awr.name}: {ex}")

    else:
        print("⚠️ No AWR HTML files found. Skipping file-based analysis.")

    # ---------------------------------------------------------------------
    # 🧮 Handle Manual Inputs (if available)
    # ---------------------------------------------------------------------
    if manual_data:
        print("---------------------------------------------------------------")
        print("🧠 Processing manual input data...")

        try:
            vcpus = float(manual_data.get("vcpus", 0))
            memory_gb = float(manual_data.get("memory_gb", 0))
            iops = float(manual_data.get("iops", 0))
            throughput = float(manual_data.get("throughput_mb_s", 0))
            db_name = "Manual_Workload"

            vm = recommend_vm_shape(cloud, vcpus, memory_gb)

            manual_result = {
                "Source": "Manual Input",
                "Cloud": cloud.upper(),
                "DB Name": db_name,
                "Memory (GB)": memory_gb,
                "Total IOPS": iops,
                "Throughput (MB/s)": throughput,
                "Estimated vCPUs": vcpus,
                "Recommended VM": vm.get("name", "N/A"),
                "VM vCPUs": vm.get("vcpus", "N/A"),
                "VM Memory (GB)": vm.get("memory", "N/A"),
                "Category": vm.get("category", "N/A"),
                "Hourly Price (USD)": vm.get("price_per_hour", "N/A"),
                "Monthly Cost (USD)": vm.get("monthly_cost", "N/A")
            }

            results.append(manual_result)
            print(f"✅ Added manual workload → {vm.get('name', 'N/A')} @ ${vm.get('price_per_hour', 'N/A')}/hr")

        except Exception as e:
            print(f"❌ Failed to process manual data: {e}")

    # ---------------------------------------------------------------------
    # 📊 Generate summary outputs
    # ---------------------------------------------------------------------
    if not results:
        print("⚠️ No data collected for summary generation.")
        return

    summary_json = OUTPUTS / "summary.json"
    summary_xlsx = OUTPUTS / "summary.xlsx"
    summary_pptx = OUTPUTS / "summary.pptx"

    # Save results
    summary_json.write_text(json.dumps(results, indent=2))
    pd.DataFrame(results).to_excel(summary_xlsx, index=False)

    # 🖼️ PowerPoint summary generation
    prs = Presentation(str(TEMPLATE_PPT))
    for r in results:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        if title:
            title.text = f"{r.get('DB Name', 'Unknown DB')} ({r.get('Cloud', '-')})"

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
        body = textbox.text_frame

        body.text = f"Source: {r.get('Source', '-')}"
        body.add_paragraph().text = f"Recommended VM: {r.get('Recommended VM', 'N/A')}"
        body.add_paragraph().text = f"Category: {r.get('Category', 'N/A')}"
        body.add_paragraph().text = f"vCPUs: {r.get('Estimated vCPUs', 'N/A')}  |  Memory: {r.get('Memory (GB)', 'N/A')} GB"
        body.add_paragraph().text = f"IOPS: {r.get('Total IOPS', 'N/A')}  |  Throughput: {r.get('Throughput (MB/s)', 'N/A')} MB/s"
        body.add_paragraph().text = f"Hourly Price: ${r.get('Hourly Price (USD)', 'N/A')}  |  Monthly Cost: ${r.get('Monthly Cost (USD)', 'N/A')}"

    prs.save(summary_pptx)

    print("===============================================================")
    print(f"✅ Summary Generated Successfully:")
    print(f"   • {summary_json}")
    print(f"   • {summary_xlsx}")
    print(f"   • {summary_pptx}")
    print("===============================================================")


# ---------------------------------------------------------------------
if __name__ == "__main__":
    import sys
    cloud = sys.argv[1] if len(sys.argv) > 1 else "azure"  # Default Azure
    analyze_all(cloud)
